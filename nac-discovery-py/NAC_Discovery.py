import json,os  
import boto3
from datetime import *
import json, logging
import pprint,re
# import requests
# from requests_aws4auth import AWS4Auth
import urllib.parse
from botocore.exceptions import ClientError
from datetime import *
import shlex,subprocess,json
from urllib.parse import unquote_plus
# import PyPDF2
from io import BytesIO
import io  
from pptx import Presentation
import fitz
from requests.auth import HTTPBasicAuth
from docx import Document
import pandas as pd
import logging
import ssl
import urllib3 
import uuid

def lambda_handler(event, context):
    logging.info('lambda_handler starts...')
    print('context.invoked_function_arn',context.invoked_function_arn)
    logging.info("Lambda function ARN:".format(context.invoked_function_arn))
    runtime_region = os.environ['AWS_REGION'] 
    context_arn=context.invoked_function_arn
    u_id=context_arn.split('-')[-1]
    logging.info('u_id'.format(u_id))

    s3 = boto3.client('s3')  
    data={}
    runtime_region = os.environ['AWS_REGION'] 
    secret_nct_nce_admin = get_secret('nasuni-labs-kendra-admin',runtime_region) 
    secret_data_internal = get_secret('nasuni-labs-internal-'+u_id,runtime_region) 
    
    share_data=read_translation_data(u_id)
    for record in event['Records']:
        logging.info(record)
        data['dest_bucket'] = record['s3']['bucket']['name']
        data['object_key'] = unquote_plus(record['s3']['object']['key'])
        data['size'] = str(record['s3']['object'].get('size', -1))
        file_name=os.path.basename(data['object_key'])
        data['file_name'] = file_name
        data['event_name'] = record['eventName']
        data['event_time'] = record['eventTime']
        data['awsRegion'] = record['awsRegion']
        data['volume_name'] = secret_data_internal['volume_name']
        obj1 = s3.get_object(Bucket=data['dest_bucket'], Key=data['object_key'])
        # data['content'] = obj1['Body'].read().decode('utf-8')
        data['extension'] = data['file_name'].split('.')[-1]
        data['root_handle'] = re.sub('[!@#$%^&*()+?=,<>/.]', '-', secret_data_internal['root_handle']).lower()
        data['source_bucket'] = secret_data_internal['discovery_source_bucket']
        logging.info("data['object_key'] = {}".format(data['object_key']))  
        logging.info("data['dest_bucket'] = {}".format(data['dest_bucket']))  
        if  data['extension'] in ['txt','csv','docx','doc','pdf','xlsx','xls','pptx','ppt']:
                
            if data['extension'] in ['csv','txt']:
                data['content'] = obj1['Body'].read().decode('utf-8')
            elif data['extension'] == 'pdf':
                file_content = obj1['Body'].read()
                text = ""
                with fitz.open(stream=file_content, filetype="pdf") as doc:
                    
                    # iterating through pdf file pages
                    for page in range(doc.page_count):
                        # fetching & appending text to text variable of each page
                        # text += page.getText()
                        text += doc.get_page_text(page) 
                    
                data['content'] = text
            elif data['extension'] in ['docx','doc']:
               fs = obj1['Body'].read()
               sentence = str(parseDocx(fs))
               logging.info('docx data {} '.format(sentence))
               data['content'] = sentence
            elif data['extension'] in ['xlsx','xls']:
                file_content = obj1['Body'].read()
                read_excel_data = io.BytesIO(file_content)
                df = pd.read_excel(read_excel_data) 
                df = df.to_string() 
                logging.info('xlsx data {}'.format(df))
                data['content'] = df 
            elif data['extension'] in ['pptx','ppt']:
                print('data[extension] elif',data['extension'])
                pptx_content = obj1['Body'].read()
                ppt = Presentation(io.BytesIO(pptx_content))
                pptx_data=''
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                pptx_data+=run.text
                logging.info("pptx data {}".format(pptx_data))
                data['content'] = pptx_data
        else:
            data['content'] =data['file_name']        
        
    print('data',data)
    ########################################################
    share_path_last_element=None
    list_after_index=None
    if share_data: 
        if share_data['path'][0]!='-' and share_data['name'][0]!='-':
            for name,path in zip(share_data['name'],share_data['path']):
                if path in data['object_key']:
                    share_path_last_element=path.split('/')[-1] 
                    logging.info('148 share_path_last_element {}'.format(share_path_last_element))
                    full_path=data['object_key']

                    full_path_with_share_name=full_path.replace(path,'/'+name)
                    logging.info('full_path_with_share_name {}'.format(full_path_with_share_name))
                    index_of_last_element=full_path_with_share_name.index(name)

                    list_after_index=full_path_with_share_name[index_of_last_element:]
                
    if secret_data_internal['web_access_appliance_address']!='not_found':
        if share_path_last_element != None:
            if share_data['name'] and share_data['path'] and share_path_last_element in data['object_key']:
                # data['access_url']='https://'+secret_data_internal['web_access_appliance_address']+'/fs/view/'+secret_data_internal['share_name']+'/'+list_after_index
                data['access_url']='https://'+secret_data_internal['web_access_appliance_address']+'/fs/view/'+list_after_index
        else:
            data['access_url']='https://'+secret_data_internal['web_access_appliance_address']+'/fs/view/'+data['volume_name']+'/'+'/'.join(data['object_key'].split('/')[3:])
    else:
        data['access_url']=secret_data_internal['web_access_appliance_address']
    
    logging.info('access_url = {}'.format(data['access_url']))
    percent_20_url=data['access_url'].replace(' ','%20')
    logging.info('After appending percent 20 url = {}'.format(percent_20_url))
    data['access_url']=percent_20_url
    logging.info('secret_data_internal = {}'.format(secret_data_internal))   
    #######################################################################
    client = boto3.client('kendra')

    #code to load the data into the kendra index.
    # index_id = '3a655609-f7ff-4d12-8521-5d2714f0db68'
    index_id=secret_nct_nce_admin['index_id']
    print('index_id',index_id)
    role_arn=secret_data_internal['discovery_lambda_role_arn']
    print('RoleArn',role_arn)  
    random_id = str(uuid.uuid4()) 
    
    response = client.batch_put_document(
        IndexId=index_id,
        Documents=[
                    {
                        'Id': random_id,
                        "Title": str(data['file_name']),
                        'S3Path': {
                            'Bucket': data['dest_bucket'],
                            'Key': data['object_key']
                        },
                        'Attributes': [
                            {
                                'Key': '_source_uri',
                                'Value': {
                                    'StringValue': data['content']+' '+data['access_url'],
                                }
                            },
                        ]
                    }
        ],
        # RoleArn= "arn:aws:iam::820345591825:role/service-role/AmazonKendra-sample-s3-role-2cbf2b4f-8f64-482d-a380-4cda5f911606"
        RoleArn= role_arn
    )

    print(response)

def read_translation_data(u_id): 
    bucket_name='nasuni-share-data-bucket-storage'    
    s3 = boto3.client('s3')
    print(bucket_name)
    
    # List all of the files in the S3 bucket
    response = s3.list_objects(Bucket=bucket_name)
    
    # Read the contents of each file in the S3 bucket
    print('response',response)
    bucket_folder_name=None
    share_data={}
    print('u_id',u_id)
    for obj in response['Contents']: 
        # Get the object key (i.e. the file name)
        key = obj['Key']
        bucket_folder_name=key
        # print('bucket_folder_name',bucket_folder_name)
        # print('key',key)  
        
        if u_id in key:
            print('found',key)
            nmc_api_filename=os.path.basename(key)
            print('nmc_api_filename',nmc_api_filename)
            if nmc_api_filename!="":
                s3.download_file(bucket_name, key, '/tmp/'+nmc_api_filename)
                
                with open('/tmp/'+nmc_api_filename, 'r') as f2:
                    if 'nmc_api_data_v_share_name' in '/tmp/'+nmc_api_filename:
                        share_data['name'] = f2.read().split(',')
                    else:
                        share_data['path'] = f2.read().split(',')
                
                    # print(data_file)
            logging.info('deleting folder from s3 bucket nasuni-share-data-bucket-storage')
            # s3.delete_object(Bucket=bucket_name, Key=key)
    print('share_data',share_data)
    logging.info(share_data)
    return share_data

def get_secret(secret_name,region_name):

    secret = ''
    session = boto3.session.Session()
    client = session.client(
        service_name='secretsmanager',
        region_name=region_name,
    )

    try:
        get_secret_value_response = client.get_secret_value(
            SecretId=secret_name
        )

    except ClientError as e:
        if e.response['Error']['Code'] == 'ResourceNotFoundException':
            print("The requested secret " + secret_name + " was not found")
        elif e.response['Error']['Code'] == 'InvalidRequestException':
            print("The request was invalid due to:", e)
        elif e.response['Error']['Code'] == 'InvalidParameterException':
            print("The request had invalid params:", e)
        elif e.response['Error']['Code'] == 'DecryptionFailure':
            print("The requested secret can't be decrypted using the provided KMS key:", e)
        elif e.response['Error']['Code'] == 'InternalServiceError':
            print("An error occurred on service side:", e)
    else:
        # Secrets Manager decrypts the secret value using the associated KMS CMK
        # Depending on whether the secret was a string or binary, only one of these fields will be populated
        if 'SecretString' in get_secret_value_response:
            secret = get_secret_value_response['SecretString']

        else:
            secret = base64.b64decode(get_secret_value_response['SecretBinary'])

    return json.loads(secret)
